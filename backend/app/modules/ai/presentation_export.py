from __future__ import annotations

import io
import re
import unicodedata
from pathlib import Path

import httpx
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from app.modules.ai.presentation_images import find_presentation_image
from app.modules.ai.schemas import GeneratedPresentationSlide, PresentationExportRequest

_PALETTES = {
    "infografico": ("EAF4FF", "1467F5", "6D4FE8"),
    "bento_pastel": ("F4F2FF", "5D67EB", "7651D8"),
    "ilustrado": ("EEF8FF", "1C70E8", "F26751"),
    "minimalista": ("FFFFFF", "111827", "4F46E5"),
    "esquema": ("F2F7FB", "1764B7", "0F8B72"),
    "alto_contraste": ("0D1830", "FFD43B", "FFE682"),
    "editorial": ("FAF7F2", "B94E3F", "31415E"),
    "gamificado": ("F1F4FF", "6045E6", "F27532"),
}


def _rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def _file_name(value: str) -> str:
    normalized = unicodedata.normalize("NFD", value)
    ascii_value = "".join(char for char in normalized if unicodedata.category(char) != "Mn")
    safe = re.sub(r"[^a-zA-Z0-9_-]+", "-", ascii_value).strip("-").lower()
    return f"{safe or 'presentacion-avendia'}.pptx"


def _asset_path(slide: GeneratedPresentationSlide) -> Path | io.BytesIO | None:
    match = re.fullmatch(r"/api/v1/ai/tools/presentation-images/([a-f0-9]{40})", slide.image_url)
    if match:
        return find_presentation_image(match.group(1))
    if slide.image_url.startswith("https://"):
        try:
            response = httpx.get(slide.image_url, timeout=20, follow_redirects=True)
            response.raise_for_status()
        except httpx.HTTPError:
            return None
        return io.BytesIO(response.content) if response.content else None
    return None


def _add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float,
    color: str,
    bold: bool = False,
    font: str = "Aptos",
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)
    return box


def _add_cover_picture(
    slide, path: Path | io.BytesIO, x: float, y: float, w: float, h: float, source_ratio: float
):
    picture = slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))
    target_ratio = w / h
    if source_ratio > target_ratio:
        crop = (1 - target_ratio / source_ratio) / 2
        picture.crop_left = crop
        picture.crop_right = crop
    elif source_ratio < target_ratio:
        crop = (1 - source_ratio / target_ratio) / 2
        picture.crop_top = crop
        picture.crop_bottom = crop
    return picture


def _add_points(
    slide, points: list[str], x: float, y: float, w: float, h: float, color: str, accent: str
):
    if not points:
        return
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    for index, point in enumerate(points[:5]):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = f"• {point}"
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(16)
        paragraph.font.bold = True
        paragraph.font.color.rgb = _rgb(color)
        paragraph.level = 0
        paragraph.space_after = Pt(10)
    return box


def _add_activity(slide, text: str, x: float, y: float, w: float, foreground: str, accent: str):
    if not text:
        return
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.07), Inches(1.05)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = _rgb(accent)
    bar.line.fill.background()
    _add_text(
        slide,
        "INTERACCIÓN EN AULA",
        x + 0.22,
        y + 0.14,
        w - 0.3,
        0.2,
        size=8.5,
        color=accent,
        bold=True,
    )
    _add_text(slide, text, x + 0.22, y + 0.42, w - 0.3, 0.52, size=10.5, color=foreground)


def build_presentation_pptx(payload: PresentationExportRequest) -> tuple[bytes, str]:
    deck = Presentation()
    deck.slide_width = Inches(13.333)
    deck.slide_height = Inches(7.5)
    blank = deck.slide_layouts[6]
    background, accent, secondary = _PALETTES[payload.visual_style]
    dark = payload.visual_style == "alto_contraste"
    base_text = "FFFFFF" if dark else "13203A"
    muted = "D7E0F0" if dark else "516078"

    for index, item in enumerate(payload.presentation.slides):
        slide = deck.slides.add_slide(blank)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = _rgb(background)
        image_path = _asset_path(item)
        source_ratio = (
            (item.image_width / item.image_height)
            if item.image_width and item.image_height
            else 16 / 9
        )
        over_photo = bool(image_path) and item.type in {"portada", "frase_destacada"}
        is_closing = item.type == "cierre"

        if image_path:
            if over_photo:
                _add_cover_picture(slide, image_path, 0, 0, 13.333, 7.5, source_ratio)
                panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(7.35), Inches(7.5))
                panel.fill.solid()
                panel.fill.fore_color.rgb = _rgb("10213F")
                panel.line.fill.background()
            elif is_closing:
                _add_cover_picture(slide, image_path, 8.08, 0.65, 4.63, 5.95, source_ratio)
            else:
                _add_cover_picture(slide, image_path, 7.65, 0, 5.683, 7.5, source_ratio)
        else:
            fallback = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(8.25), 0, Inches(5.083), Inches(7.5)
            )
            fallback.fill.solid()
            fallback.fill.fore_color.rgb = _rgb(accent)
            fallback.line.fill.background()
            _add_text(
                slide,
                f"{index + 1:02d}",
                8.55,
                2.1,
                4.25,
                2.1,
                size=80,
                color="FFFFFF",
                bold=True,
                align=PP_ALIGN.CENTER,
            )

        foreground = "FFFFFF" if over_photo else base_text
        foreground_muted = "E6EDFA" if over_photo else muted
        label_color = "BCD0FF" if over_photo else secondary
        copy_x = 1.05 if item.type == "frase_destacada" else 0.68
        copy_w = (
            8.25
            if item.type == "frase_destacada"
            else 6.35
            if image_path and not over_photo
            else 6.5
        )
        _add_text(
            slide,
            f"{index + 1:02d} · {payload.curricular_area.upper()}",
            copy_x,
            0.5,
            copy_w,
            0.25,
            size=9,
            color=label_color,
            bold=True,
        )
        title_y = 1.25 if item.type == "portada" else 0.92
        title_size = 32 if item.type == "portada" else 27 if item.type == "frase_destacada" else 24
        _add_text(
            slide,
            item.title,
            copy_x,
            title_y,
            copy_w,
            1.45 if item.type == "portada" else 1.0,
            size=title_size,
            color=foreground,
            bold=True,
            font="Aptos Display",
            valign=MSO_ANCHOR.MIDDLE,
        )
        if item.subtitle:
            _add_text(
                slide,
                item.subtitle,
                copy_x,
                2.86 if item.type == "portada" else 1.92,
                copy_w,
                0.6,
                size=13,
                color=foreground_muted,
            )
        # Una portada solo comunica el tema. El desarrollo, cita y actividad se distribuyen
        # en las siguientes diapositivas para que la exportación conserve la misma jerarquía
        # visual que el editor.
        if item.type not in {"portada", "frase_destacada"}:
            _add_points(
                slide,
                item.key_points,
                copy_x,
                3.55 if item.type == "portada" else 2.55,
                copy_w,
                2.25 if item.interactive_activity else 3.2,
                foreground,
                accent,
            )
        if item.type != "portada" and item.highlighted_quote:
            _add_text(
                slide,
                item.highlighted_quote,
                copy_x,
                3.02 if item.type == "frase_destacada" else 4.65,
                copy_w,
                1.55 if item.type == "frase_destacada" else 0.9,
                size=25 if item.type == "frase_destacada" else 17,
                color="FFFFFF" if item.type == "frase_destacada" else accent,
                bold=True,
                valign=MSO_ANCHOR.MIDDLE,
            )
        if item.type != "portada":
            _add_activity(
                slide,
                item.interactive_activity,
                copy_x,
                5.62,
                min(copy_w, 6.5),
                foreground,
                "FFFFFF" if over_photo else accent,
            )
        _add_text(
            slide,
            f"{payload.teacher_name} · {payload.institution}",
            0.68,
            7.08,
            6.4,
            0.15,
            size=7,
            color=foreground_muted,
        )
        if image_path and item.image_attribution:
            _add_text(
                slide,
                item.image_attribution,
                8.0,
                7.08,
                4.55,
                0.15,
                size=5.5,
                color="FFFFFF" if over_photo else "68758A",
                align=PP_ALIGN.RIGHT,
            )
        _add_text(
            slide,
            f"{index + 1}/{len(payload.presentation.slides)}",
            12.65,
            7.07,
            0.35,
            0.16,
            size=7,
            color="FFFFFF" if over_photo else muted,
            bold=True,
            align=PP_ALIGN.RIGHT,
        )
        notes = slide.notes_slide.notes_text_frame
        notes.text = item.speaker_notes
        if item.image_attribution:
            notes.text += f"\n\nImagen: {item.image_attribution}"
            if item.image_source_url:
                notes.text += f"\nFuente: {item.image_source_url}"

    output = io.BytesIO()
    deck.save(output)
    return output.getvalue(), _file_name(payload.presentation.presentation_title)
