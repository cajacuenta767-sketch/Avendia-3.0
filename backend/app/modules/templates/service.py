import re
from io import BytesIO
from pathlib import Path
from xml.sax.saxutils import escape

from docx import Document as WordDocument
from docx.enum.text import WD_BREAK
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader, PdfWriter
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.units import mm
from reportlab.platypus import PageBreak, Paragraph, SimpleDocTemplate, Spacer

from app.modules.ai.schemas import GeneratedWorkflowArtifact
from app.modules.templates.model import InstitutionalTemplate

MIME_TYPES = {
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".pdf": "application/pdf",
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}


def safe_stem(name: str) -> str:
    stem = Path(name).stem
    cleaned = re.sub(r"[^a-zA-Z0-9áéíóúñÁÉÍÓÚÑ _-]", "", stem).strip()
    return cleaned or "formato-institucional"


def artifact_text(artifact: GeneratedWorkflowArtifact) -> str:
    parts = [artifact.document_title, artifact.executive_summary]
    for section in artifact.sections:
        parts.extend([section.title, section.narrative, *section.key_points])
    parts.extend(["Recomendaciones para revisión docente", *artifact.teacher_recommendations])
    return "\n\n".join(parts)


def replace_word_placeholders(document: WordDocument, artifact: GeneratedWorkflowArtifact) -> bool:
    replacements = {
        "{{titulo}}": artifact.document_title,
        "{{resumen}}": artifact.executive_summary,
        "{{contenido}}": artifact_text(artifact),
    }
    replaced = False
    for paragraph in document.paragraphs:
        text = paragraph.text
        next_text = text
        for placeholder, value in replacements.items():
            next_text = next_text.replace(placeholder, value)
        if next_text != text:
            paragraph.text = next_text
            replaced = True
    for table in document.tables:
        for row in table.rows:
            for cell in row.cells:
                text = cell.text
                next_text = text
                for placeholder, value in replacements.items():
                    next_text = next_text.replace(placeholder, value)
                if next_text != text:
                    cell.text = next_text
                    replaced = True
    return replaced


def render_docx(template: InstitutionalTemplate, artifact: GeneratedWorkflowArtifact) -> bytes:
    document = WordDocument(BytesIO(template.content))
    content_replaced = replace_word_placeholders(document, artifact)
    if not content_replaced:
        paragraph = document.add_paragraph()
        paragraph.add_run().add_break(WD_BREAK.PAGE)
        document.add_heading(artifact.document_title, level=1)
        document.add_paragraph(artifact.executive_summary)
        for section in artifact.sections:
            document.add_heading(section.title, level=2)
            document.add_paragraph(section.narrative)
            for point in section.key_points:
                document.add_paragraph(point, style="List Bullet")
        document.add_heading("Recomendaciones para revisión docente", level=2)
        for recommendation in artifact.teacher_recommendations:
            document.add_paragraph(recommendation, style="List Bullet")
    output = BytesIO()
    document.save(output)
    return output.getvalue()


def render_pdf(template: InstitutionalTemplate, artifact: GeneratedWorkflowArtifact) -> bytes:
    generated = BytesIO()
    styles = getSampleStyleSheet()
    story = [
        Paragraph(escape(artifact.document_title), styles["Title"]),
        Spacer(1, 6 * mm),
        Paragraph(escape(artifact.executive_summary), styles["BodyText"]),
    ]
    for section in artifact.sections:
        story.extend(
            [
                Spacer(1, 5 * mm),
                Paragraph(escape(section.title), styles["Heading2"]),
                Paragraph(escape(section.narrative), styles["BodyText"]),
            ]
        )
        story.extend(
            Paragraph(f"• {escape(point)}", styles["BodyText"]) for point in section.key_points
        )
    story.extend(
        [PageBreak(), Paragraph("Recomendaciones para revisión docente", styles["Heading2"])]
    )
    story.extend(
        Paragraph(f"• {escape(item)}", styles["BodyText"])
        for item in artifact.teacher_recommendations
    )
    SimpleDocTemplate(
        generated,
        pagesize=A4,
        leftMargin=18 * mm,
        rightMargin=18 * mm,
        topMargin=18 * mm,
        bottomMargin=18 * mm,
    ).build(story)
    writer = PdfWriter()
    for page in PdfReader(BytesIO(template.content)).pages:
        writer.add_page(page)
    for page in PdfReader(BytesIO(generated.getvalue())).pages:
        writer.add_page(page)
    output = BytesIO()
    writer.write(output)
    return output.getvalue()


def render_xlsx(template: InstitutionalTemplate, artifact: GeneratedWorkflowArtifact) -> bytes:
    workbook = load_workbook(BytesIO(template.content))
    if "Contenido Avendia" in workbook.sheetnames:
        del workbook["Contenido Avendia"]
    sheet = workbook.create_sheet("Contenido Avendia")
    rows: list[tuple[str, str]] = [
        ("Título", artifact.document_title),
        ("Resumen", artifact.executive_summary),
    ]
    for section in artifact.sections:
        rows.append((section.title, section.narrative))
        rows.extend(("Punto clave", point) for point in section.key_points)
    rows.extend(
        ("Recomendación", recommendation) for recommendation in artifact.teacher_recommendations
    )
    for row_index, (label, content) in enumerate(rows, start=1):
        sheet.cell(row=row_index, column=1, value=label)
        sheet.cell(row=row_index, column=2, value=content)
    sheet.column_dimensions["A"].width = 30
    sheet.column_dimensions["B"].width = 110
    output = BytesIO()
    workbook.save(output)
    return output.getvalue()


def add_presentation_slide(presentation: Presentation, title: str, body: str) -> None:
    layout = (
        presentation.slide_layouts[1]
        if len(presentation.slide_layouts) > 1
        else presentation.slide_layouts[0]
    )
    slide = presentation.slides.add_slide(layout)
    if slide.shapes.title:
        slide.shapes.title.text = title
    placeholders = [
        shape
        for shape in slide.placeholders
        if shape != slide.shapes.title and hasattr(shape, "text_frame")
    ]
    if placeholders:
        placeholders[0].text = body[:4500]


def render_pptx(template: InstitutionalTemplate, artifact: GeneratedWorkflowArtifact) -> bytes:
    presentation = Presentation(BytesIO(template.content))
    add_presentation_slide(presentation, artifact.document_title, artifact.executive_summary)
    for section in artifact.sections:
        body = "\n".join([section.narrative, *[f"• {point}" for point in section.key_points]])
        add_presentation_slide(presentation, section.title, body)
    output = BytesIO()
    presentation.save(output)
    return output.getvalue()


def render_template(
    template: InstitutionalTemplate, artifact: GeneratedWorkflowArtifact
) -> tuple[bytes, str, str]:
    renderers = {
        ".docx": render_docx,
        ".pdf": render_pdf,
        ".xlsx": render_xlsx,
        ".pptx": render_pptx,
    }
    renderer = renderers[template.extension]
    content = renderer(template, artifact)
    filename = f"{safe_stem(template.name)}-avendia{template.extension}"
    return content, MIME_TYPES[template.extension], filename
