from io import BytesIO

import pytest
from docx import Document
from httpx import ASGITransport, AsyncClient
from openpyxl import Workbook, load_workbook
from pptx import Presentation
from pypdf import PdfReader
from reportlab.pdfgen import canvas

from app.main import app


def registration(email: str) -> dict[str, object]:
    return {
        "email": email,
        "full_name": "María Gómez",
        "password": "a-secure-password",
        "dre": "DRE Lima Metropolitana",
        "ugel": "UGEL 03",
        "school_name": "I.E. José María Arguedas",
        "director_name": "Elena Torres",
        "education_modality": "EBR",
        "education_level": "Secundaria",
        "grade": "3° de Secundaria",
        "section": "A",
        "curricular_area": "Comunicación",
        "school_year": 2026,
    }


async def authenticated_headers(client: AsyncClient, email: str) -> dict[str, str]:
    payload = registration(email)
    assert (await client.post("/api/v1/auth/register", json=payload)).status_code == 201
    login = await client.post(
        "/api/v1/auth/login",
        json={"email": email, "password": payload["password"]},
    )
    return {"Authorization": f"Bearer {login.json()['access_token']}"}


def template_docx() -> bytes:
    document = Document()
    document.add_heading("Formato institucional", level=1)
    document.add_paragraph("{{titulo}}")
    document.add_paragraph("{{resumen}}")
    document.add_paragraph("{{contenido}}")
    output = BytesIO()
    document.save(output)
    return output.getvalue()


def template_pdf() -> bytes:
    output = BytesIO()
    document = canvas.Canvas(output)
    document.drawString(72, 760, "Formato institucional")
    document.save()
    return output.getvalue()


def template_xlsx() -> bytes:
    workbook = Workbook()
    workbook.active.title = "Carátula"
    workbook.active["A1"] = "Formato institucional"
    output = BytesIO()
    workbook.save(output)
    return output.getvalue()


def template_pptx() -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    slide.shapes.title.text = "Formato institucional"
    output = BytesIO()
    presentation.save(output)
    return output.getvalue()


def artifact() -> dict[str, object]:
    return {
        "document_title": "Plan lector institucional",
        "executive_summary": (
            "Propuesta institucional de lectura que organiza actividades y evidencias "
            "para el periodo escolar vigente."
        ),
        "sections": [
            {
                "title": "Propósito",
                "narrative": (
                    "Fortalecer la comprensión lectora mediante experiencias progresivas "
                    "y contextualizadas."
                ),
                "key_points": ["Lectura semanal", "Evidencias de comprensión"],
            }
        ],
        "teacher_recommendations": [
            "Revisar las fechas con el calendario institucional.",
            "Validar los criterios con el equipo docente.",
        ],
        "activity": None,
    }


@pytest.mark.asyncio
async def test_template_upload_is_owner_scoped_and_renders_docx() -> None:
    async with AsyncClient(transport=ASGITransport(app=app), base_url="http://test") as client:
        owner = await authenticated_headers(client, "template-owner@example.edu")
        other = await authenticated_headers(client, "template-other@example.edu")

        uploaded = await client.post(
            "/api/v1/templates",
            headers=owner,
            files={
                "file": (
                    "formato-colegio.docx",
                    template_docx(),
                    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                )
            },
            data={"make_default": "true"},
        )
        assert uploaded.status_code == 201
        template_id = uploaded.json()["id"]
        assert uploaded.json()["is_default"] is True

        listing = await client.get("/api/v1/templates", headers=owner)
        assert listing.status_code == 200
        assert [item["name"] for item in listing.json()] == ["formato-colegio.docx"]

        forbidden = await client.get(f"/api/v1/templates/{template_id}/download", headers=other)
        assert forbidden.status_code == 404

        rendered = await client.post(
            f"/api/v1/templates/{template_id}/render",
            headers=owner,
            json={"artifact": artifact(), "document_type": "plan-curricular-anual"},
        )
        assert rendered.status_code == 200
        rendered_document = Document(BytesIO(rendered.content))
        rendered_text = "\n".join(paragraph.text for paragraph in rendered_document.paragraphs)
        assert "Plan lector institucional" in rendered_text
        assert "{{titulo}}" not in rendered_text

        deleted = await client.delete(f"/api/v1/templates/{template_id}", headers=owner)
        assert deleted.status_code == 204
        assert (await client.get("/api/v1/templates", headers=owner)).json() == []


@pytest.mark.asyncio
async def test_template_render_supports_pdf_xlsx_and_pptx() -> None:
    cases = [
        ("formato.pdf", "application/pdf", template_pdf()),
        (
            "formato.xlsx",
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            template_xlsx(),
        ),
        (
            "formato.pptx",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            template_pptx(),
        ),
    ]
    async with AsyncClient(transport=ASGITransport(app=app), base_url="http://test") as client:
        headers = await authenticated_headers(client, "template-formats@example.edu")
        rendered_by_extension: dict[str, bytes] = {}
        for filename, content_type, content in cases:
            uploaded = await client.post(
                "/api/v1/templates",
                headers=headers,
                files={"file": (filename, content, content_type)},
            )
            assert uploaded.status_code == 201
            rendered = await client.post(
                f"/api/v1/templates/{uploaded.json()['id']}/render",
                headers=headers,
                json={"artifact": artifact(), "document_type": "sesion-aprendizaje"},
            )
            assert rendered.status_code == 200
            rendered_by_extension[filename.rsplit(".", 1)[1]] = rendered.content

        assert len(PdfReader(BytesIO(rendered_by_extension["pdf"])).pages) >= 2
        workbook = load_workbook(BytesIO(rendered_by_extension["xlsx"]))
        assert "Contenido Avendia" in workbook.sheetnames
        presentation = Presentation(BytesIO(rendered_by_extension["pptx"]))
        assert len(presentation.slides) >= 2
